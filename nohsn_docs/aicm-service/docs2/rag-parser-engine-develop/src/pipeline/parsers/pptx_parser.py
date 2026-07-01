"""PPTX 파서 -- python-pptx 기반 슬라이드별 텍스트/표/이미지 추출.

전략
-----
1. 각 슬라이드를 개별 PageContent 로 변환 (slide index = page_number)
2. 텍스트: shape.has_text_frame → 단락별 TextBlock (타이틀/서브타이틀은 heading 감지)
3. 표: shape.has_table → TableContent (첫 행 = 헤더, markdown 변환)
4. 이미지: Picture shape → ImageContent (임시 디렉터리에 저장, alt_text 활용)
5. 발표자 노트: slide.notes_slide → 메타데이터로 보존
"""

from __future__ import annotations

import asyncio
import os
import tempfile
from typing import Any

from src.common.logging import get_logger
from src.pipeline.models.parse_result import (
    ImageContent,
    PageContent,
    ParseResult,
    TableContent,
    TextBlock,
)
from src.pipeline.parsers.base import BaseParser

log = get_logger(__name__)

# 타이틀 placeholder idx (0 = 제목, 1 = 부제목/본문 타이틀)
_TITLE_PLACEHOLDER_IDXS = {0, 1}


class PPTXParser(BaseParser):
    """PPTX 프레젠테이션 파서."""

    async def parse(self) -> ParseResult:
        """PPTX 파일을 파싱하여 슬라이드별 구조화 결과를 반환한다."""
        return await asyncio.to_thread(self._parse_sync)

    # ------------------------------------------------------------------
    # 동기 내부
    # ------------------------------------------------------------------
    def _parse_sync(self) -> ParseResult:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:
            raise RuntimeError(
                "python-pptx 패키지가 필요합니다: pip install python-pptx"
            ) from exc

        prs = Presentation(self.file_path)

        pages: list[PageContent] = []
        all_tables: list[TableContent] = []
        all_images: list[ImageContent] = []
        all_texts: list[str] = []
        slide_notes: dict[int, str] = {}

        # 이미지 저장용 임시 디렉터리
        img_dir = tempfile.mkdtemp(prefix="aicm_pptx_img_")

        for slide_idx, slide in enumerate(prs.slides):
            page_number = slide_idx + 1
            text_blocks: list[TextBlock] = []
            page_tables: list[TableContent] = []
            page_images: list[ImageContent] = []
            slide_texts: list[str] = []
            running_offset = 0
            table_idx = 0
            image_idx = 0

            for shape in slide.shapes:
                # --- 표 추출 ---
                if shape.has_table:
                    tc = self._extract_table(
                        shape.table, page_number, len(all_tables) + table_idx
                    )
                    if tc:
                        page_tables.append(tc)
                        table_idx += 1
                    continue

                # --- 이미지 추출 ---
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    ic = self._extract_image(
                        shape, page_number, len(all_images) + image_idx, img_dir
                    )
                    if ic:
                        page_images.append(ic)
                        image_idx += 1
                    continue

                # --- 텍스트 추출 ---
                if shape.has_text_frame:
                    heading_level = self._detect_heading(shape)
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        para_text = para.text.strip()
                        if not para_text:
                            continue

                        # 첫 번째 단락만 heading level 적용 (타이틀 shape인 경우)
                        block_heading = heading_level if para_idx == 0 else None

                        tb = TextBlock(
                            text=para_text,
                            start_char_offset=running_offset,
                            end_char_offset=running_offset + len(para_text),
                            paragraph_index=len(text_blocks),
                            heading_level=block_heading,
                        )
                        text_blocks.append(tb)
                        slide_texts.append(para_text)
                        running_offset += len(para_text) + 1

            # --- 발표자 노트 ---
            notes_text = self._extract_notes(slide)
            if notes_text:
                slide_notes[page_number] = notes_text

            page_text = "\n".join(slide_texts)
            all_texts.append(page_text)

            pages.append(
                PageContent(
                    page_number=page_number,
                    text=page_text,
                    text_blocks=text_blocks,
                    tables=page_tables,
                    images=page_images,
                    layout_type="single",
                )
            )

            all_tables.extend(page_tables)
            all_images.extend(page_images)

        full_text = "\n\n".join(all_texts)

        # 메타데이터
        metadata = self._extract_metadata(prs)
        if slide_notes:
            metadata["slide_notes"] = slide_notes

        log.info(
            "pptx_parse_complete",
            slide_count=len(pages),
            total_blocks=sum(len(p.text_blocks) for p in pages),
            table_count=len(all_tables),
            image_count=len(all_images),
            file=self.file_path,
        )

        return ParseResult(
            raw_text=full_text,
            pages=pages,
            tables=all_tables,
            images=all_images,
            metadata=metadata,
            source_file_path=self.file_path,
        )

    # ------------------------------------------------------------------
    # 헤딩 감지
    # ------------------------------------------------------------------
    @staticmethod
    def _detect_heading(shape: Any) -> int | None:
        """슬라이드 타이틀/서브타이틀 shape 여부로 헤딩 레벨을 판정한다."""
        if not shape.is_placeholder:
            return None
        try:
            ph_idx = shape.placeholder_format.idx
        except Exception:
            return None
        if ph_idx == 0:
            return 1  # 슬라이드 제목
        if ph_idx == 1:
            return 2  # 부제목
        return None

    # ------------------------------------------------------------------
    # 표 추출
    # ------------------------------------------------------------------
    @staticmethod
    def _extract_table(
        table: Any, page_number: int, table_index: int
    ) -> TableContent | None:
        """python-pptx Table 객체에서 TableContent 를 추출한다."""
        all_rows: list[list[str]] = []
        for row in table.rows:
            all_rows.append([cell.text.replace("\n", " ").strip() for cell in row.cells])

        if not all_rows:
            return None

        headers = all_rows[0]
        body = all_rows[1:] if len(all_rows) > 1 else []

        # markdown 변환
        lines = [
            "| " + " | ".join(headers) + " |",
            "| " + " | ".join("---" for _ in headers) + " |",
        ]
        for row in body:
            padded = row + [""] * max(0, len(headers) - len(row))
            lines.append("| " + " | ".join(padded[: len(headers)]) + " |")
        markdown = "\n".join(lines)

        return TableContent(
            page_number=page_number,
            table_index=table_index,
            headers=headers,
            rows=body,
            markdown=markdown,
            confidence=0.95,
        )

    # ------------------------------------------------------------------
    # 이미지 추출
    # ------------------------------------------------------------------
    @staticmethod
    def _extract_image(
        shape: Any, page_number: int, image_index: int, img_dir: str
    ) -> ImageContent | None:
        """Picture shape 에서 이미지를 추출하여 임시 파일로 저장한다."""
        try:
            image = shape.image
            content_type = image.content_type
            # 확장자 결정
            ext_map = {
                "image/png": ".png",
                "image/jpeg": ".jpg",
                "image/gif": ".gif",
                "image/bmp": ".bmp",
                "image/tiff": ".tiff",
                "image/webp": ".webp",
                "image/svg+xml": ".svg",
            }
            ext = ext_map.get(content_type, ".png")
            filename = f"slide{page_number}_img{image_index}{ext}"
            img_path = os.path.join(img_dir, filename)

            with open(img_path, "wb") as f:
                f.write(image.blob)

            # alt_text 추출
            description = None
            try:
                alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                if alt_text.strip():
                    description = alt_text.strip()
            except Exception:
                pass

            # name 속성에서도 시도
            if not description:
                try:
                    name = shape.name
                    if name and not name.startswith("Picture"):
                        description = name
                except Exception:
                    pass

            return ImageContent(
                page_number=page_number,
                image_index=image_index,
                image_path=img_path,
                description=description,
            )
        except Exception:
            log.warning(
                "pptx_image_extraction_failed",
                page_number=page_number,
                image_index=image_index,
                exc_info=True,
            )
            return None

    # ------------------------------------------------------------------
    # 발표자 노트
    # ------------------------------------------------------------------
    @staticmethod
    def _extract_notes(slide: Any) -> str:
        """슬라이드의 발표자 노트를 추출한다."""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                if text_frame and text_frame.text.strip():
                    return text_frame.text.strip()
        except Exception:
            pass
        return ""

    # ------------------------------------------------------------------
    # 메타데이터
    # ------------------------------------------------------------------
    @staticmethod
    def _extract_metadata(prs: Any) -> dict:
        """프레젠테이션 전체 메타데이터를 추출한다."""
        metadata: dict[str, Any] = {}

        # 코어 프로퍼티
        try:
            cp = prs.core_properties
            metadata["title"] = cp.title or ""
            metadata["author"] = cp.author or ""
            metadata["created"] = str(cp.created) if cp.created else ""
            metadata["modified"] = str(cp.modified) if cp.modified else ""
            metadata["subject"] = cp.subject or ""
        except Exception:
            pass

        # 슬라이드 수
        metadata["slide_count"] = len(prs.slides)

        # 슬라이드 크기
        try:
            metadata["slide_width_emu"] = prs.slide_width
            metadata["slide_height_emu"] = prs.slide_height
            # EMU → cm 변환 (1 EMU = 1/914400 inch = 1/360000 cm)
            metadata["slide_width_cm"] = round(prs.slide_width / 360000, 2)
            metadata["slide_height_cm"] = round(prs.slide_height / 360000, 2)
        except Exception:
            pass

        # 슬라이드 레이아웃 이름 목록
        try:
            layout_names: list[str] = []
            for slide in prs.slides:
                layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
                layout_names.append(layout_name)
            metadata["slide_layouts"] = layout_names
        except Exception:
            pass

        return metadata
